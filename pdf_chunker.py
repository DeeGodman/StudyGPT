import os
import uuid
import json
import fitz  # PyMuPDF
import docx
from pptx import Presentation
from textwrap import wrap

INPUT_DIR = "SLIDES"
OUTPUT_FILE = "chunks_output.json"
CHUNK_SIZE = 500  # token-ish chunk size

def extract_text_from_pdf(path):
    doc = fitz.open(path)
    return "\n".join([page.get_text() for page in doc])

def extract_text_from_docx(path):
    doc = docx.Document(path)
    return "\n".join([para.text for para in doc.paragraphs if para.text.strip()])

def extract_text_from_pptx(path):
    prs = Presentation(path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
    return "\n".join(texts)

def chunk_text(text, chunk_size=CHUNK_SIZE):
    words = text.split()
    return wrap(" ".join(words), chunk_size * 5)

def get_file_type_extractor(ext):
    return {
        ".pdf": extract_text_from_pdf,
        ".docx": extract_text_from_docx,
        ".pptx": extract_text_from_pptx
    }.get(ext.lower())

def extract_metadata(file_path):
    parts = file_path.split(os.sep)
    level = next((p for p in parts if "level_" in p), "unknown")
    semester = next((p for p in parts if "sem" in p), "unknown")
    return level, semester

def process_all_files(input_dir):
    all_chunks = []
    for root, _, files in os.walk(input_dir):
        for filename in files:
            ext = os.path.splitext(filename)[1]
            extractor = get_file_type_extractor(ext)
            if not extractor:
                continue
            full_path = os.path.join(root, filename)
            print(f"Processing {full_path}...")
            try:
                text = extractor(full_path)
            except Exception as e:
                print(f"Failed to extract from {filename}: {e}")
                continue
            chunks = chunk_text(text)
            course = filename.split("_")[0]
            week = filename.split("_")[1] if "_" in filename else "unknown"
            level, semester = extract_metadata(full_path)
            for i, chunk in enumerate(chunks):
                all_chunks.append({
                    "id": str(uuid.uuid4()),
                    "course": course,
                    "week": week,
                    "level": level,
                    "semester": semester,
                    "file": filename,
                    "file_type": ext.lower().replace(".", ""),
                    "chunk_index": i,
                    "text": chunk.strip()
                })
    return all_chunks

if __name__ == "__main__":
    chunks = process_all_files(INPUT_DIR)
    with open(OUTPUT_FILE, "w", encoding="utf-8") as f:
        json.dump(chunks, f, indent=2)
    print(f"✅ Done. Extracted {len(chunks)} chunks into {OUTPUT_FILE}")
